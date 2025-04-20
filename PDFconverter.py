from pptx import Presentation
from pptx.util import Inches, Pt
import os
import fitz  # PyMuPDF
from tkinter import Tk, Button, Label, PhotoImage, Toplevel, StringVar
from tkinter.filedialog import askopenfilename, asksaveasfilename
from tkinter.ttk import Progressbar

def pdf_to_pptx(pdf_path, pptx_path, progress_var, progress_bar, status_label):
    # Verificar si el archivo PDF existe
    if not os.path.exists(pdf_path):
        status_label.config(text=f"El archivo {pdf_path} no existe.")
        return

    # Crear una presentación PPTX
    presentation = Presentation()

    # Abrir el archivo PDF
    pdf_document = fitz.open(pdf_path)

    # Ajustar el tamaño de la diapositiva al tamaño del PDF
    first_page = pdf_document[0]
    pdf_width, pdf_height = first_page.rect.width, first_page.rect.height
    presentation.slide_width = Pt(pdf_width)
    presentation.slide_height = Pt(pdf_height)

    total_pages = len(pdf_document)
    for page_number in range(total_pages):
        # Actualizar la barra de progreso
        progress_var.set((page_number + 1) / total_pages * 100)
        progress_bar.update()

        # Obtener la página del PDF
        page = pdf_document[page_number]

        # Renderizar la página como imagen
        pix = page.get_pixmap(dpi=150)  # DPI ajustable
        image_path = f"page_{page_number + 1}.png"
        pix.save(image_path)

        # Crear una diapositiva en el PPTX
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Layout vacío

        # Agregar la imagen a la diapositiva
        left = Inches(0)
        top = Inches(0)
        slide.shapes.add_picture(image_path, left, top, Pt(pdf_width), Pt(pdf_height))

        # Eliminar la imagen temporal
        os.remove(image_path)

    # Guardar el archivo PPTX
    presentation.save(pptx_path)
    status_label.config(text=f"Conversión completada. Archivo guardado en: {pptx_path}")

def center_window(window, width, height):
    """Centrar una ventana en la pantalla."""
    screen_width = window.winfo_screenwidth()
    screen_height = window.winfo_screenheight()
    x = (screen_width // 2) - (width // 2)
    y = (screen_height // 2) - (height // 2)
    window.geometry(f"{width}x{height}+{x}+{y}")

def open_file_selection():
    # Deshabilitar los widgets de la ventana principal
    start_button.config(state="disabled")

    # Abrir ventana para seleccionar el archivo PDF
    pdf_path = askopenfilename(
        title="Selecciona el archivo PDF",
        filetypes=[("Archivos PDF", "*.pdf")]
    )

    if not pdf_path:
        print("No se seleccionó ningún archivo PDF.")
        start_button.config(state="normal")  # Habilitar el botón nuevamente
    else:
        # Abrir ventana para seleccionar la ruta de guardado del PPTX
        pptx_path = asksaveasfilename(
            title="Guardar como",
            defaultextension=".pptx",
            filetypes=[("Archivos PPTX", "*.pptx")]
        )

        if not pptx_path:
            print("No se seleccionó una ruta para guardar el archivo PPTX.")
            start_button.config(state="normal")  # Habilitar el botón nuevamente
        else:
            # Crear una ventana de progreso
            progress_window = Toplevel(root)
            progress_window.title("Progreso de conversión")
            center_window(progress_window, 400, 200)  # Centrar la ventana de progreso

            progress_label = Label(progress_window, text="Convirtiendo PDF a PPTX...")
            progress_label.pack(pady=10)

            progress_var = StringVar()
            progress_bar = Progressbar(progress_window, orient="horizontal", length=300, mode="determinate", variable=progress_var)
            progress_bar.pack(pady=10)

            status_label = Label(progress_window, text="")
            status_label.pack(pady=10)

            # Función para habilitar el botón de cerrar al finalizar
            def on_conversion_complete():
                close_button = Button(progress_window, text="Cerrar", command=progress_window.destroy,relief="raised",  # Estilo tridimensional
                bg="#4CAF50",  # Color de fondo (verde)
                fg="white",  # Color del texto
                font=("Arial", 12, "bold"))
                close_button.pack(pady=10)
                start_button.config(state="normal")  # Habilitar el botón principal nuevamente

            # Realizar la conversión
            pdf_to_pptx(pdf_path, pptx_path, progress_var, progress_bar, status_label)

            # Mostrar el botón de cerrar al finalizar
            progress_window.after(100, on_conversion_complete)

if __name__ == "__main__":
    # Crear la ventana principal
    root = Tk()
    root.title("PDF a PPTX Converter")
    center_window(root, 700, 500)  # Centrar la ventana principal

    # Cargar el logo y reducir su tamaño
    logo_image = PhotoImage(file="logo.png")
    resized_logo = logo_image.subsample(3, 3)  # Reducir el tamaño del logo
    logo_label = Label(root, image=resized_logo)
    logo_label.pack(pady=10)

    # Etiqueta de bienvenida
    label = Label(root, text="Bienvenido al convertidor PDF a PPTX")
    label.pack(pady=10)

    # Botón para iniciar la selección de archivos
    start_button = Button(
        root,
        text="Seleccionar PDF",
        command=open_file_selection,
        relief="raised",  # Estilo tridimensional
        bg="#4CAF50",  # Color de fondo (verde)
        fg="white",  # Color del texto
        font=("Arial", 12, "bold")  # Fuente personalizada
    )
    start_button.pack(pady=20)

    # Iniciar el bucle principal de la ventana
    root.mainloop()